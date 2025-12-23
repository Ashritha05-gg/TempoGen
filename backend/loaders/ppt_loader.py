# # loaders/ppt_loader.py

# from pptx import Presentation
# import os
# from typing import List, Dict

# def extract_ppt_text(file_path: str) -> List[Dict]:
#     prs = Presentation(file_path)
#     output = []
#     filename = os.path.basename(file_path)

#     for i, slide in enumerate(prs.slides):
#         slide_text = []

#         for shape in slide.shapes:
#             if hasattr(shape, "text"):
#                 slide_text.append(shape.text)

#         if slide_text:
#             output.append({
#                 "text": "\n".join(slide_text),
#                 "page": f"slide-{i+1}",
#                 "source": filename
#             })

#     return output


from pptx import Presentation
import os
from typing import List, Dict

def extract_ppt_text(file_path: str) -> List[Dict]:
    prs = Presentation(file_path)
    output = []
    filename = os.path.basename(file_path)

    for i, slide in enumerate(prs.slides):
        slide_text = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())

        if slide_text:
            text = f"""
Slide {i+1} Content:
{chr(10).join(slide_text)}
""".strip()

            output.append({
                "text": text,
                "page": f"slide-{i+1}",
                "source": filename
            })

    return output
